# pages/analisis.py
import streamlit as st
import pandas as pd
import numpy as np
import pickle
import statsmodels.api as sm
import shutil
import os
import plotly.express as px
import io
import base64
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import tempfile
import time
import plotly.io as pio

st.set_page_config(page_title="RevFlux", page_icon="Logo.png", layout="wide")

# Redirect otomatis jika tidak ada dataset aktif

# ============================================================
# Ambil active_dataset dari session_state atau file
# ============================================================
active_dataset = st.session_state.get("active_dataset", None)
if not active_dataset and os.path.exists("active_dataset.txt"):
    try:
        with open("active_dataset.txt", "r") as f:
            active_dataset = f.read().strip()
    except Exception:
        active_dataset = None
        
if not active_dataset:
    st.warning("⚠️ Tidak ada dataset aktif. Mengarahkan ke halaman utama...")
    time.sleep(1)
    st.switch_page("index.py")

# ============================================================
# Fungsi preprocess (asli, robust)
# ============================================================
def preprocess_period_column(df):
    """Mempersiapkan kolom 'Periode' agar konsisten sebagai datetime (awal bulan)."""
    df = df.copy()
    if "Periode" not in df.columns:
        raise ValueError("Kolom 'Periode' tidak ditemukan.")
    # Parse ke datetime
    df["Periode"] = pd.to_datetime(df["Periode"], errors="coerce", infer_datetime_format=True)
    # Jika masih ada NaT, coba format-format umum
    if df["Periode"].isna().any():
        for fmt in ("%d/%m/%Y", "%Y/%m/%d", "%Y-%m", "%Y-%m-%d", "%d-%m-%Y", "%m/%Y", "%Y"):
            parsed = pd.to_datetime(df["Periode"].astype(str), format=fmt, errors="coerce")
            df["Periode"] = df["Periode"].fillna(parsed)
            if df["Periode"].notna().all():
                break

    df = df.dropna(subset=["Periode"]).copy()
    # samakan ke awal bulan
    df["Periode"] = df["Periode"].dt.to_period("M").dt.to_timestamp()

    if "Pemasukan" in df.columns:
        df["Pemasukan"] = pd.to_numeric(df["Pemasukan"], errors="coerce")
        df = df.dropna(subset=["Pemasukan"])

    df = df.sort_values("Periode").reset_index(drop=True)
    return df

# ============================================================
# Load existing model (jika ada)
# ============================================================
sarima_model = None
if active_dataset and os.path.exists(f"{active_dataset}_model.pkl"):
    try:
        with open(f"{active_dataset}_model.pkl", "rb") as f:
            sarima_model = pickle.load(f)
    except Exception:
        sarima_model = None

# ============================================================
# Styling & navbar (sama seperti original)
# ============================================================
logo_path = Path("Logo.png")
if logo_path.exists():
    with open(logo_path, "rb") as f:
        logo_base64 = base64.b64encode(f.read()).decode()
    logo_html = f'<img src="data:image/png;base64,{logo_base64}" style="width:28px;height:28px;margin-right:10px;">'
else:
    logo_html = "<div style='width:28px;height:28px;background:#ccc;border-radius:50%;margin-right:10px;'></div>"

navbar_css = """
<style>
header[data-testid="stHeader"], footer, div[data-testid="stToolbar"] {display:none;}
.navbar {
  position: fixed; top: 0; left: 0; width: 100%; height: 60px;
  display: flex; justify-content: flex-start; align-items: center;
  padding: 0 24px; z-index: 9999; border-radius: 0 0 10px 10px;
  background-color: #f8f9fa; color: #333;
  box-shadow: 0 2px 8px rgba(0,0,0,0.1);
}
.navbar span {font-weight:700;font-size:20px;color:#8e44ad;}
.block-container {padding-top:80px !important;}
.divider {border:none;height:2px;margin:30px 0 20px 0;background:linear-gradient(to right, #f8f9fa);}
div.stDownloadButton>button, div.stButton>button {
    border-radius:8px !important; font-weight:600 !important;
}
div.stButton>button:has(span:contains('Train Model')),
div.stButton>button:has(span:contains('Retrain Model')),
div.stDownloadButton>button {
    background-color:#8e44ad !important; color:white !important; border:none !important;
}
div.stDownloadButton>button:hover, div.stButton>button:hover {opacity:0.9 !important;}
div.stButton>button:has(span:contains('Reset Model')) {
    background-color:#f39c12 !important; color:white !important;
}
div.stButton>button:has(span:contains('Kembalikan Model')) {
    background-color:#3498db !important; color:white !important;
}
</style>
"""
st.markdown(navbar_css, unsafe_allow_html=True)
navbar_html = f"<div class='navbar'>{logo_html}<span>RevFlux</span></div>"
st.markdown(navbar_html, unsafe_allow_html=True)

# ============================================================
# Pembuka
# ============================================================
st.markdown(
    """ 
    <div style='text-align: center;'> 
        <h1 style='font-size:40px; color:#8e44ad; font-weight:800;'>RevFlux</h1> 
        <p style='font-size:22px; font-weight:600; color:#333;'>
            Insight & Prediksi Penjualan
        </p> 
        <hr style='margin-top:20px; margin-bottom:20px; border: 1px solid #8e44ad;'/> 
    </div> 
    """,
    unsafe_allow_html=True
)

# ============================================================
# Pastikan ada active_dataset & data file
# ============================================================

data_filename = f"{active_dataset}_data.csv"
if not os.path.exists(data_filename):
    st.error(f"File data '{data_filename}' tidak ditemukan. Silakan pastikan Anda telah mengupload data lewat halaman utama.")
    st.stop()

# ============================================================
# Load data awal dan tampilkan ringkasan
# ============================================================
sales_data = pd.read_csv(data_filename)
try:
    sales_data = preprocess_period_column(sales_data)
except Exception as e:
    st.error(f"Gagal memproses kolom Periode/Pemasukan: {e}")
    st.stop()

last_period = sales_data["Periode"].iloc[-1].strftime("%B %Y")
st.info(f"📅 Data terakhir: **{last_period}**")
st.success(f"📊 Dataset aktif: {active_dataset} ({len(sales_data)} baris)")
st.dataframe(sales_data.tail(8))

st.markdown("<hr class='divider'/>", unsafe_allow_html=True)

# ============================================================
# Bagian: Tambah Data Baru — Hanya Upload CSV/XLSX (Opsi A)
# ============================================================
st.subheader("➕ Tambah Data Baru — Upload CSV/XLSX (Periode, Pemasukan)")

# jika baru saja selesai save, tampilkan info sekali
if st.session_state.get("data_saved_ok", False):
    st.info("📌 Dataset telah berhasil diperbarui.")
    st.session_state["data_saved_ok"] = False
else:
    additional_file = st.file_uploader(
        "Upload file CSV/XLSX yang berisi kolom 'Periode' & 'Pemasukan'",
        type=["csv", "xlsx"],
        key="additional_upload"
    )

    if additional_file:
        try:
            # read
            if additional_file.name.lower().endswith(".csv"):
                add_df = pd.read_csv(additional_file)
            else:
                add_df = pd.read_excel(additional_file)

            # validasi header
            if "Periode" not in add_df.columns or "Pemasukan" not in add_df.columns:
                st.error("❌ File tambahan harus memiliki kolom 'Periode' dan 'Pemasukan'. Pastikan header tepat.")
            else:
                # preprocess keseuaian periode
                add_df = preprocess_period_column(add_df)

                # merge
                old = pd.read_csv(data_filename)
                old = preprocess_period_column(old)
                combined = pd.concat([old, add_df], ignore_index=True)
                combined = combined.drop_duplicates(subset=["Periode"], keep="last") \
                                   .sort_values("Periode") \
                                   .reset_index(drop=True)

                # simpan
                combined.to_csv(data_filename, index=False, date_format="%Y-%m-%d")

                # tandai sebagai sukses
                st.session_state["data_saved_ok"] = True

                # clean upload session
                for k in ["additional_upload"]:
                    if k in st.session_state:
                        del st.session_state[k]

                st.rerun()

        except Exception as e:
            st.error(f"Gagal menambahkan file: {e}")

st.markdown("<hr class='divider'/>", unsafe_allow_html=True)


# ============================================================
# Tombol Train / Retrain Model
# ============================================================
st.subheader("🚀 Latih Model")
if active_dataset and os.path.exists(f"{active_dataset}_data.csv"):
    data_filename = f"{active_dataset}_data.csv"
    model_filename = f"{active_dataset}_model.pkl"

    sales_data = pd.read_csv(data_filename)
    sales_data = preprocess_period_column(sales_data)
    last_period = sales_data["Periode"].iloc[-1].strftime("%B %Y")
    
    st.info(f"📅 Data terakhir (setelah update): **{last_period}**")

    train_button_label = "🚀 Train Model Baru" if sarima_model is None else "🔁 Retrain Model"
    if st.button(train_button_label):
        with st.spinner("Sedang melatih model... 🧠"):
            try:
                if len(sales_data) < 24:
                    st.warning("⚠️ Jumlah data disarankan minimal 24 bulan untuk hasil optimal.")

                y = sales_data["Pemasukan"]
                y_log = np.log1p(y)
                model = sm.tsa.statespace.SARIMAX(
                    y_log, order=(1,1,1), seasonal_order=(1,1,1,12),
                    enforce_stationarity=False, enforce_invertibility=False
                ).fit(disp=False)
                with open(model_filename, "wb") as f:
                    pickle.dump(model, f)
                sarima_model = model
                st.success(f"✅ Model untuk '{active_dataset}' berhasil dilatih dengan data hingga {last_period}!")
            except Exception as e:
                st.error(f"Gagal melatih model: {e}")
else:
    st.info("⚠️ Tidak ada data untuk dilatih.")

st.markdown("<hr class='divider'/>", unsafe_allow_html=True)

# ============================================================
# Prediksi & Visualisasi
# ============================================================
if sarima_model is not None and active_dataset and os.path.exists(f"{active_dataset}_data.csv"):
    st.subheader("📈 Prediksi & Visualisasi")
    st.markdown("<hr class='divider'/>", unsafe_allow_html=True)

    n_periods = st.slider("Pilih jumlah bulan ke depan untuk prediksi:", 1, 24, 6)

    # --- Persiapan Data untuk Visualisasi ---
    hist_data = pd.read_csv(f"{active_dataset}_data.csv")
    hist_data = preprocess_period_column(hist_data)
    hist_data["Tipe"] = "Penjualan Sebelumnya"

    forecast_res = sarima_model.get_forecast(steps=n_periods)
    forecast_mean = np.expm1(forecast_res.predicted_mean)
    conf_int_exp = np.expm1(forecast_res.conf_int())

    forecast_df = pd.DataFrame({
        "Periode": pd.date_range(hist_data["Periode"].iloc[-1] + pd.DateOffset(months=1), periods=n_periods, freq="MS"),
        "Pemasukan": forecast_mean,
        "Tipe": "Prediksi"
    })
    
    combined_vis = pd.concat([hist_data, forecast_df], ignore_index=True)
    
    st.write("Tabel Hasil Prediksi:")
    display_df = forecast_df.copy()
    display_df["Periode"] = display_df["Periode"].dt.strftime("%B %Y")
    display_df["Pemasukan (Rp)"] = display_df["Pemasukan"].apply(lambda x: f"Rp {x:,.0f}".replace(",", "."))
    st.dataframe(display_df[["Periode", "Pemasukan (Rp)"]])

    # --- Buat dan Tampilkan Semua Grafik ---
    st.write("---")
    fig_line = px.line(combined_vis, x="Periode", y="Pemasukan", color="Tipe", markers=True, 
                       title=f"📈 Prediksi Pemasukan — Dataset: {active_dataset}",
                       color_discrete_map={"Penjualan Sebelumnya": "#3498db", "Prediksi": "#e74c3c"})
    fig_line.update_layout(legend_title_text="Jenis Data", yaxis_title="Pemasukan (Rp)", xaxis_title="Periode")
    st.plotly_chart(fig_line, use_container_width=True)

    st.write("---")
    fig_ci = px.line(hist_data, x="Periode", y="Pemasukan", title="🎯 Prediksi Pemasukan dengan Rentang Keyakinan (Confidence Interval)")
    fig_ci.data[0].name = 'Penjualan Sebelumnya'
    fig_ci.data[0].showlegend = True
    fig_ci.add_scatter(x=forecast_df["Periode"], y=forecast_df["Pemasukan"], mode="lines", name="Prediksi", line=dict(color="#e74c3c"))
    fig_ci.add_scatter(x=forecast_df["Periode"], y=conf_int_exp.iloc[:, 1], mode="lines", line=dict(dash="dash", color="green"), name="Batas Atas CI")
    fig_ci.add_scatter(x=forecast_df["Periode"], y=conf_int_exp.iloc[:, 0], mode="lines", line=dict(dash="dash", color="yellow"), name="Batas Bawah CI")
    fig_ci.update_layout(yaxis_title="Pemasukan (Rp)", xaxis_title="Periode", legend_title_text="Keterangan")
    st.plotly_chart(fig_ci, use_container_width=True)

    st.write("---")
    tail_periods = st.slider("Tampilkan N bulan terakhir pada Bar Chart:", 6, 36, 12)
    bar_data = combined_vis.tail(tail_periods)
    
    fig_bar = px.bar(bar_data, x="Periode", y="Pemasukan", color="Tipe",
                     barmode="group", title=f"📊 Perbandingan Pemasukan Bulanan (Penjualan Sebelumnya vs Prediksi) — {tail_periods} Bulan Terakhir",
                     labels={"Pemasukan": "Pemasukan (Rp)"},
                     color_discrete_map={"Penjualan Sebelumnya": "#3498db", "Prediksi": "#e74c3c"})
    fig_bar.update_layout(xaxis_title="Periode", yaxis_title="Pemasukan (Rp)", legend_title_text="Jenis Data")
    st.plotly_chart(fig_bar, use_container_width=True)

    st.markdown("<hr class='divider'/>", unsafe_allow_html=True)


    # ============================================================
    # Export PowerPoint
    # ============================================================
    st.subheader("💾 Export Visualisasi")
    progress_bar = None

    if st.button("📤 Export Visualisasi ke PowerPoint"):
        progress_bar = st.progress(0)
        try:
            prs = Presentation()

            # SLIDE 1 — Judul + Info Dataset
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Laporan Prediksi Pemasukan"

            try:
                slide.placeholders[1].text = (
                    f"Dataset: {active_dataset}\n"
                    f"Periode Data Aktual: {hist_data['Periode'].min().strftime('%B %Y')} - "
                    f"{hist_data['Periode'].max().strftime('%B %Y')}\n"
                    f"Periode Prediksi: {n_periods} bulan ke depan"
                )
            except Exception:
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
                tf = txBox.text_frame
                tf.text = (
                    f"Dataset: {active_dataset}\n"
                    f"Periode Data Aktual: {hist_data['Periode'].min().strftime('%B %Y')} - "
                    f"{hist_data['Periode'].max().strftime('%B %Y')}\n"
                    f"Periode Prediksi: {n_periods} bulan ke depan"
                )

            progress_bar.progress(20)
            time.sleep(0.3)

            # SLIDE 2 — Ringkasan Prediksi
            slide_summary = prs.slides.add_slide(prs.slide_layouts[1])
            title_shape = slide_summary.shapes.title
            if title_shape:
                title_shape.text = "Ringkasan Prediksi"
            else:
                tbox = slide_summary.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
                tbox.text = "Ringkasan Prediksi"

            content_box = slide_summary.shapes.add_textbox(Inches(1), Inches(2), Inches(8.5), Inches(4))
            content_frame = content_box.text_frame
            content_frame.text = (
                "Berikut adalah ringkasan hasil analisis dan prediksi pemasukan:\n"
                f"- Dataset aktif: {active_dataset}\n"
                f"- Jumlah periode historis: {len(hist_data)} bulan\n"
                f"- Periode data aktual: {hist_data['Periode'].min().strftime('%B %Y')} - "
                f"{hist_data['Periode'].max().strftime('%B %Y')}\n"
                f"- Prediksi untuk {n_periods} bulan ke depan.\n\n"
                "Grafik berikut akan memperlihatkan perbandingan antara data aktual, hasil prediksi, "
                "dan rentang keyakinan model."
            )

            progress_bar.progress(40)
            time.sleep(0.3)

            # SLIDE 3-5 — Grafik Visualisasi
            figs = [
                (fig_line, "Grafik Aktual vs Prediksi"),
                (fig_ci, "Grafik Rentang Keyakinan Prediksi"),
                (fig_bar, "Grafik Perbandingan Bulanan")
            ]

            for i, (fig, title) in enumerate(figs):
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                title_shape = slide.shapes.title
                if title_shape:
                    title_shape.text = title
                else:
                    tbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
                    tbox.text = title

                img_bytes = pio.to_image(fig, format="png", width=960, height=540)
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
                    tmp_img.write(img_bytes)
                    slide.shapes.add_picture(tmp_img.name, Inches(0.5), Inches(1.2), width=Inches(9))

                progress_bar.progress(50 + int((i + 1) * 15))
                time.sleep(0.2)

            # Simpan dan tampilkan tombol download
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmpfile:
                prs.save(tmpfile.name)

            progress_bar.progress(100)
            time.sleep(0.5)

            with open(tmpfile.name, "rb") as f:
                st.download_button(
                    label="⬇️ Download Laporan PowerPoint",
                    data=f.read(),
                    file_name=f"laporan_prediksi_{active_dataset}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

            os.remove(tmpfile.name)
            st.success("✅ PowerPoint berhasil dibuat!")

        except Exception as e:
            st.error(f"Gagal membuat PowerPoint: {e}")

        finally:
            if progress_bar is not None:
                time.sleep(0.5)
                progress_bar.empty()

else:
    st.markdown("<hr class='divider'/>", unsafe_allow_html=True)
    st.info("💡 Fitur **Export Visualisasi ke PowerPoint** akan tersedia setelah Anda melatih model dan memiliki data aktif.")

st.markdown("<hr class='divider'/>", unsafe_allow_html=True)

# ============================================================
# ⚙️ Pengaturan Sistem & Model (gabungkan Reset, Restore, Inisialisasi)
# ============================================================
st.subheader("⚙️ Pengaturan Sistem & Model")
st.markdown("<hr class='divider'/>", unsafe_allow_html=True)

col_a, col_b = st.columns(2)

with col_a:
    if st.button("🧹 Reset Model & Data"):
        try:
            if active_dataset:
                data_file = f"{active_dataset}_data.csv"
                model_file = f"{active_dataset}_model.pkl"
                backup_data = f"{active_dataset}_data_backup.csv"
                backup_model = f"{active_dataset}_model_backup.pkl"

                # Simpan backup dulu
                if os.path.exists(model_file):
                    shutil.copy(model_file, backup_model)
                if os.path.exists(data_file):
                    shutil.copy(data_file, backup_data)

                # Hapus file aktif
                for f in [model_file, data_file, "active_dataset.txt"]:
                    if os.path.exists(f):
                        os.remove(f)

                # ============== ADD THIS PART ==============
                # Tulis flag reset agar INDEX bisa menampilkan tombol restore
                with open("reset_flag.txt","w") as f:
                    f.write("1")

                st.session_state["data_saved_ok"] = False
                # ===========================================

                # Hapus session state terkait
                if "active_dataset" in st.session_state:
                    del st.session_state["active_dataset"]
                if "data_saved_ok" in st.session_state:
                    del st.session_state["data_saved_ok"]

                st.success(f"Model dan data '{active_dataset}' berhasil direset.")
                st.info("Halaman akan direfresh...")
                time.sleep(1.5)
                st.switch_page("index.py")   # <<-- gunakan ini agar langsung ke index

            else:
                st.warning("⚠️ Tidak ada dataset aktif untuk direset.")
        except Exception as e:
            st.error(f"Gagal mereset: {e}")


# Inisialisasi Ulang Sistem (hapus semua file sementara & session)
if st.button("🔄 Inisialisasi Ulang Sistem (Hapus semua data & model)"):
    try:
        protected_files = {"requirements.txt", "packages.txt", "runtime.txt", "Logo.png"}
        for file in os.listdir():
            if file in protected_files:
                continue
            if file.endswith((".pkl", ".csv", ".txt", ".backup.pkl", ".backup.csv")):
                try:
                    os.remove(file)
                except Exception:
                    pass

        # Hapus tmp files terkait ppt/png
        tmp_dir = tempfile.gettempdir()
        for f in os.listdir(tmp_dir):
            if f.startswith("tmp") and f.endswith((".pptx", ".png")):
                try:
                    os.remove(os.path.join(tmp_dir, f))
                except Exception:
                    pass

        # Clear session_state
        for key in list(st.session_state.keys()):
            del st.session_state[key]

        st.success("✅ Semua data, model, dan file sementara telah dihapus.")
        st.info("Memuat ulang aplikasi...")
        time.sleep(2)
        st.switch_page("index.py")
    except Exception as e:
        st.error(f"Gagal melakukan inisialisasi ulang: {e}")
